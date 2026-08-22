"""Least-privilege loaders for files explicitly selected in Agent Chat UI."""

from __future__ import annotations

import csv
import io
import json
import os
import posixpath
import uuid
import zipfile
from collections.abc import Iterator
from email import policy
from email.parser import BytesParser
from pathlib import Path, PurePosixPath

from bs4 import BeautifulSoup
from defusedxml import ElementTree
from docx import Document as WordDocument
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from odf import teletype
from odf.opendocument import load as load_odf
from odf.table import Table, TableCell, TableRow
from odf.text import H, P
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

from src.epub_attachments import EpubAttachmentError, extract_epub

MAX_ATTACHMENT_BYTES = 25 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 5_000
MAX_EXPANDED_BYTES = 150 * 1024 * 1024
MAX_SEGMENT_CHARACTERS = 100_000
MAX_TOTAL_CHARACTERS = 300_000
OCR_UPLOAD_DIR = Path(
    os.getenv(
        "OCR_UPLOAD_DIR",
        str(Path(__file__).resolve().parents[2] / "data" / "ocr" / "uploads"),
    )
).resolve()


def preserve_ocr_upload(data: bytes, filename: str) -> dict:
    """Persist an explicitly uploaded document and return a non-guessable reference."""
    if len(data) > MAX_ATTACHMENT_BYTES:
        raise DocumentAttachmentError("The selected file exceeds the 25 MB limit")
    safe_name = PurePosixPath(filename).name or "attachment"
    # Validate bytes and format before making the file available to OCR.
    SelectedFileLoader(data, safe_name)
    OCR_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    path = (OCR_UPLOAD_DIR / f"{uuid.uuid4().hex}-{safe_name}").resolve()
    path.write_bytes(data)
    return {"reference": f"upload:{path.name}", "path": str(path), "filename": safe_name}


TEXT_EXTENSIONS = {
    ".txt",
    ".text",
    ".md",
    ".markdown",
    ".mdx",
    ".rst",
    ".adoc",
    ".asciidoc",
    ".tex",
    ".bib",
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
    ".ndjson",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".conf",
    ".env.example",
    ".log",
    ".sql",
    ".graphql",
    ".gql",
    ".py",
    ".pyi",
    ".js",
    ".jsx",
    ".mjs",
    ".cjs",
    ".ts",
    ".tsx",
    ".java",
    ".kt",
    ".kts",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".swift",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".ps1",
    ".bat",
    ".cmd",
    ".c",
    ".h",
    ".cc",
    ".cpp",
    ".hpp",
    ".cs",
    ".dart",
    ".lua",
    ".r",
    ".jl",
    ".ex",
    ".exs",
    ".erl",
    ".hrl",
    ".clj",
    ".cljs",
    ".vue",
    ".svelte",
    ".vtt",
    ".srt",
}

DOCUMENT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".odt",
    ".ods",
    ".odp",
    ".rtf",
    ".html",
    ".htm",
    ".xhtml",
    ".xml",
    ".eml",
    ".ipynb",
    ".epub",
}

REJECTED_EXTENSIONS = {
    ".lcpl": "Readium LCP licenses must be opened by a compliant reading system",
    ".docm": "Macro-enabled Office documents are not accepted",
    ".xlsm": "Macro-enabled Office documents are not accepted",
    ".pptm": "Macro-enabled Office documents are not accepted",
    ".zip": "Archives must be unpacked by the human before selecting a file",
    ".rar": "Archives must be unpacked by the human before selecting a file",
    ".7z": "Archives must be unpacked by the human before selecting a file",
    ".tar": "Archives must be unpacked by the human before selecting a file",
    ".gz": "Archives must be unpacked by the human before selecting a file",
}


class DocumentAttachmentError(ValueError):
    """A concise, recoverable selected-file validation error."""


def supported_extensions() -> list[str]:
    return sorted(TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS)


def _extension(filename: str) -> str:
    lower = PurePosixPath(filename).name.lower()
    if lower.endswith(".env.example"):
        return ".env.example"
    return Path(lower).suffix


def _decode_text(data: bytes) -> str:
    if b"\x00" in data[:8192]:
        raise DocumentAttachmentError("The selected file appears to be binary")
    for encoding in ("utf-8-sig", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentAttachmentError("The selected text encoding is not supported")


def _validate_zip_container(data: bytes) -> None:
    try:
        archive = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise DocumentAttachmentError(
            "The selected document container is malformed"
        ) from exc
    with archive:
        entries = archive.infolist()
        if len(entries) > MAX_ARCHIVE_ENTRIES:
            raise DocumentAttachmentError(
                "The document contains too many embedded entries"
            )
        if sum(entry.file_size for entry in entries) > MAX_EXPANDED_BYTES:
            raise DocumentAttachmentError(
                "The document expanded size exceeds the safety limit"
            )
        for entry in entries:
            path = entry.filename.replace("\\", "/")
            parts = PurePosixPath(path).parts
            if path.startswith("/") or ".." in parts:
                raise DocumentAttachmentError(
                    "The document contains an unsafe embedded path"
                )
            if entry.flag_bits & 0x1:
                raise DocumentAttachmentError(
                    "Encrypted document containers are not supported"
                )
            if posixpath.basename(path).lower() == "vbaproject.bin":
                raise DocumentAttachmentError(
                    "Documents containing macros are not accepted"
                )


def _document(text: str, filename: str, segment: str, **metadata) -> Document:
    return Document(
        page_content=text[:MAX_SEGMENT_CHARACTERS],
        metadata={
            "filename": PurePosixPath(filename).name,
            "segment": segment,
            **metadata,
        },
    )


class SelectedFileLoader(BaseLoader):
    """Load only supplied bytes; paths and directories are intentionally unsupported."""

    def __init__(self, data: bytes, filename: str):
        if len(data) > MAX_ATTACHMENT_BYTES:
            raise DocumentAttachmentError("The selected file exceeds the 25 MB limit")
        self.data = data
        self.filename = PurePosixPath(filename).name
        self.extension = _extension(self.filename)
        if self.extension in REJECTED_EXTENSIONS:
            raise DocumentAttachmentError(REJECTED_EXTENSIONS[self.extension])
        if self.extension not in TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS:
            raise DocumentAttachmentError("This file format is not supported")
        if self.extension in {
            ".docx",
            ".xlsx",
            ".pptx",
            ".odt",
            ".ods",
            ".odp",
            ".epub",
        }:
            _validate_zip_container(data)

    def lazy_load(self) -> Iterator[Document]:
        extension = self.extension
        if extension in TEXT_EXTENSIONS:
            yield _document(_decode_text(self.data), self.filename, "document")
        elif extension == ".pdf":
            yield from self._pdf()
        elif extension == ".docx":
            yield from self._docx()
        elif extension == ".xlsx":
            yield from self._xlsx()
        elif extension == ".pptx":
            yield from self._pptx()
        elif extension in {".odt", ".ods", ".odp"}:
            yield from self._open_document()
        elif extension == ".rtf":
            yield _document(
                rtf_to_text(_decode_text(self.data)), self.filename, "document"
            )
        elif extension in {".html", ".htm", ".xhtml"}:
            yield from self._html()
        elif extension == ".xml":
            yield from self._xml()
        elif extension == ".eml":
            yield from self._email()
        elif extension == ".ipynb":
            yield from self._notebook()
        elif extension == ".epub":
            yield from self._epub()

    def _pdf(self) -> Iterator[Document]:
        try:
            reader = PdfReader(io.BytesIO(self.data))
            if reader.is_encrypted:
                raise DocumentAttachmentError("Encrypted PDFs are not supported")
            for index, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    yield _document(text, self.filename, f"page-{index}", page=index)
        except DocumentAttachmentError:
            raise
        except Exception as exc:
            raise DocumentAttachmentError("The PDF could not be read safely") from exc

    def _docx(self) -> Iterator[Document]:
        try:
            document = WordDocument(io.BytesIO(self.data))
            paragraphs = [
                paragraph.text for paragraph in document.paragraphs if paragraph.text
            ]
            for table_index, table in enumerate(document.tables, start=1):
                paragraphs.append(f"\nTable {table_index}")
                paragraphs.extend(
                    " | ".join(cell.text for cell in row.cells) for row in table.rows
                )
            yield _document("\n".join(paragraphs), self.filename, "document")
        except Exception as exc:
            raise DocumentAttachmentError(
                "The DOCX document could not be read safely"
            ) from exc

    def _xlsx(self) -> Iterator[Document]:
        try:
            workbook = load_workbook(
                io.BytesIO(self.data), read_only=True, data_only=True, keep_links=False
            )
            for sheet in workbook.worksheets:
                output = io.StringIO()
                writer = csv.writer(output)
                for row in sheet.iter_rows(values_only=True):
                    writer.writerow("" if value is None else value for value in row)
                yield _document(
                    output.getvalue(),
                    self.filename,
                    f"sheet-{sheet.title}",
                    sheet=sheet.title,
                )
            workbook.close()
        except Exception as exc:
            raise DocumentAttachmentError(
                "The XLSX workbook could not be read safely"
            ) from exc

    def _pptx(self) -> Iterator[Document]:
        try:
            presentation = Presentation(io.BytesIO(self.data))
            for index, slide in enumerate(presentation.slides, start=1):
                text = "\n".join(
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
                if text:
                    yield _document(text, self.filename, f"slide-{index}", slide=index)
        except Exception as exc:
            raise DocumentAttachmentError(
                "The PPTX presentation could not be read safely"
            ) from exc

    def _open_document(self) -> Iterator[Document]:
        try:
            document = load_odf(io.BytesIO(self.data))
            if self.extension == ".ods":
                for table_index, table in enumerate(
                    document.getElementsByType(Table), start=1
                ):
                    rows = []
                    for row in table.getElementsByType(TableRow):
                        rows.append(
                            " | ".join(
                                teletype.extractText(cell)
                                for cell in row.getElementsByType(TableCell)
                            )
                        )
                    yield _document(
                        "\n".join(rows),
                        self.filename,
                        f"sheet-{table_index}",
                        sheet=table_index,
                    )
            else:
                elements = document.getElementsByType(H) + document.getElementsByType(P)
                text = "\n".join(teletype.extractText(element) for element in elements)
                yield _document(text, self.filename, "document")
        except Exception as exc:
            raise DocumentAttachmentError(
                "The OpenDocument file could not be read safely"
            ) from exc

    def _html(self) -> Iterator[Document]:
        soup = BeautifulSoup(_decode_text(self.data), "html.parser")
        for element in soup(["script", "style", "noscript", "template"]):
            element.decompose()
        yield _document(soup.get_text("\n", strip=True), self.filename, "document")

    def _xml(self) -> Iterator[Document]:
        try:
            root = ElementTree.fromstring(self.data)
        except Exception as exc:
            raise DocumentAttachmentError(
                "The XML document could not be read safely"
            ) from exc
        yield _document(
            "\n".join(text.strip() for text in root.itertext() if text.strip()),
            self.filename,
            "document",
        )

    def _email(self) -> Iterator[Document]:
        message = BytesParser(policy=policy.default).parsebytes(self.data)
        headers = [
            f"{name}: {message.get(name, '')}"
            for name in ("Subject", "From", "To", "Date")
            if message.get(name)
        ]
        bodies = []
        parts = message.walk() if message.is_multipart() else (message,)
        for part in parts:
            if part.get_content_type() == "text/plain" and not part.get_filename():
                bodies.append(part.get_content())
        yield _document("\n".join([*headers, "", *bodies]), self.filename, "message")

    def _notebook(self) -> Iterator[Document]:
        try:
            notebook = json.loads(_decode_text(self.data))
        except json.JSONDecodeError as exc:
            raise DocumentAttachmentError("The notebook JSON is malformed") from exc
        for index, cell in enumerate(notebook.get("cells", []), start=1):
            source = cell.get("source", [])
            text = "".join(source) if isinstance(source, list) else str(source)
            if text:
                yield _document(
                    text,
                    self.filename,
                    f"cell-{index}",
                    cell=index,
                    cell_type=cell.get("cell_type", "unknown"),
                )

    def _epub(self) -> Iterator[Document]:
        try:
            result = extract_epub(self.data, self.filename)
        except EpubAttachmentError as exc:
            raise DocumentAttachmentError(str(exc)) from exc
        if not result["text"]:
            raise DocumentAttachmentError(
                "This EPUB has no extractable text; image-only publications need explicit OCR"
            )
        yield _document(
            result["text"],
            self.filename,
            "publication",
            title=result["title"],
            author=result["author"],
            chapters=result["chapters"],
            content_profile=result["content_profile"],
            truncated=result["truncated"],
        )


def load_selected_document(data: bytes, filename: str) -> dict:
    documents = SelectedFileLoader(data, filename).load()
    if not documents:
        raise DocumentAttachmentError("No readable text was found in the selected file")

    parts: list[str] = []
    segments: list[dict] = []
    remaining = MAX_TOTAL_CHARACTERS
    truncated = False
    for document in documents:
        text = document.page_content.strip()
        if not text:
            continue
        accepted = text[:remaining]
        metadata = dict(document.metadata)
        segment = str(metadata.pop("segment", f"segment-{len(segments) + 1}"))
        parts.append(f"## {segment}\n\n{accepted}")
        segments.append({"id": segment, "characters": len(text), **metadata})
        remaining -= len(accepted)
        if remaining <= 0:
            truncated = len(accepted) < len(text) or len(segments) < len(documents)
            break

    if not parts:
        raise DocumentAttachmentError("No readable text was found in the selected file")
    return {
        "filename": PurePosixPath(filename).name,
        "format": _extension(filename).removeprefix("."),
        "text": "\n\n".join(parts),
        "segments": segments,
        "truncated": truncated,
    }
