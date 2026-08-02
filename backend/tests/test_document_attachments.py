import io
import json

import pytest
from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation

from src.document_attachments import (
    DocumentAttachmentError,
    load_selected_document,
    supported_extensions,
)


def test_plain_text_and_source_code_are_normalized_without_a_path():
    result = load_selected_document(
        b"def answer():\n    return 42\n", "/private/code.py"
    )

    assert result["filename"] == "code.py"
    assert result["format"] == "py"
    assert "return 42" in result["text"]
    assert "/private" not in result["text"]


def test_docx_tables_are_extracted():
    document = WordDocument()
    document.add_paragraph("Project overview")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Component"
    table.rows[0].cells[1].text = "Jasper"
    payload = io.BytesIO()
    document.save(payload)

    result = load_selected_document(payload.getvalue(), "overview.docx")

    assert "Project overview" in result["text"]
    assert "Component | Jasper" in result["text"]


def test_xlsx_sheets_are_extracted_as_rows():
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Architecture"
    sheet.append(["Node", "Purpose"])
    sheet.append(["Jasper", "Supervisor"])
    payload = io.BytesIO()
    workbook.save(payload)

    result = load_selected_document(payload.getvalue(), "map.xlsx")

    assert "Node,Purpose" in result["text"]
    assert result["segments"][0]["sheet"] == "Architecture"


def test_pptx_slides_are_extracted():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Request flow"
    slide.placeholders[1].text = "UI to LangGraph"
    payload = io.BytesIO()
    presentation.save(payload)

    result = load_selected_document(payload.getvalue(), "flow.pptx")

    assert "Request flow" in result["text"]
    assert "UI to LangGraph" in result["text"]


def test_html_removes_executable_content():
    result = load_selected_document(
        b"<html><body><h1>Useful</h1><script>secret()</script></body></html>",
        "page.html",
    )

    assert "Useful" in result["text"]
    assert "secret()" not in result["text"]


def test_email_and_notebook_are_segmented():
    email = load_selected_document(
        b"Subject: Notes\nFrom: human@example.test\n\nSession findings", "notes.eml"
    )
    notebook = load_selected_document(
        json.dumps(
            {
                "cells": [
                    {"cell_type": "markdown", "source": ["# Architecture"]},
                    {"cell_type": "code", "source": ["print('hello')"]},
                ]
            }
        ).encode(),
        "analysis.ipynb",
    )

    assert "Session findings" in email["text"]
    assert len(notebook["segments"]) == 2


@pytest.mark.parametrize("filename", ["book.lcpl", "macro.docm", "bundle.zip"])
def test_unsafe_or_authority_expanding_formats_are_rejected(filename):
    with pytest.raises(DocumentAttachmentError):
        load_selected_document(b"not accepted", filename)


def test_xml_entities_are_rejected_and_supported_formats_are_discoverable():
    with pytest.raises(DocumentAttachmentError, match="safely"):
        load_selected_document(
            b'<!DOCTYPE x [<!ENTITY secret SYSTEM "file:///etc/passwd">]><x>&secret;</x>',
            "unsafe.xml",
        )

    extensions = supported_extensions()
    assert {".pdf", ".docx", ".xlsx", ".pptx", ".epub", ".py"} <= set(extensions)
    assert ".lcpl" not in extensions
